# make_ppt.py  —  generates SmartCampus_DSA_Presentation.pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────
BG       = RGBColor(0x1E, 0x1E, 0x2E)   # slide background
BG2      = RGBColor(0x18, 0x18, 0x25)   # code block bg
ACCENT   = RGBColor(0x0F, 0x6E, 0x56)   # teal green
PURPLE   = RGBColor(0xCB, 0xA6, 0xF7)   # purple
BLUE     = RGBColor(0x89, 0xB4, 0xFA)   # blue
YELLOW   = RGBColor(0xF9, 0xE2, 0xAF)   # yellow
RED      = RGBColor(0xF3, 0x8B, 0xA8)   # red/pink
GREEN    = RGBColor(0xA6, 0xE3, 0xA1)   # green
WHITE    = RGBColor(0xCD, 0xD6, 0xF4)   # off-white text
GRAY     = RGBColor(0x45, 0x47, 0x5A)   # border gray
DARK     = RGBColor(0x11, 0x11, 0x1B)   # darkest

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ──────────────────────────────────────────────────
def add_slide():
    sl = prs.slides.add_slide(BLANK)
    fill_bg(sl)
    return sl

def fill_bg(sl, color=BG):
    bg = sl.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(sl, x, y, w, h, bg=None, border=None):
    from pptx.util import Inches
    shp = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.line.fill.background()
    if bg:
        shp.fill.solid(); shp.fill.fore_color.rgb = bg
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp

def txt(sl, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def title_bar(sl, title, subtitle=None):
    """Green left-border title bar."""
    box(sl, 0.3, 0.15, 0.08, 0.7, bg=ACCENT)
    txt(sl, title, 0.55, 0.12, 11, 0.55, size=28, bold=True, color=PURPLE)
    if subtitle:
        txt(sl, subtitle, 0.55, 0.65, 11, 0.4, size=14, color=WHITE)

def code_block(sl, code, x, y, w, h, title=None):
    """Dark code block with monospace text."""
    box(sl, x, y, w, h, bg=BG2, border=GRAY)
    if title:
        txt(sl, title, x+0.1, y+0.05, w-0.2, 0.28, size=11,
            color=ACCENT, bold=True)
        yoff = 0.32
    else:
        yoff = 0.1
    txb = sl.shapes.add_textbox(Inches(x+0.15), Inches(y+yoff),
                                 Inches(w-0.3), Inches(h-yoff-0.1))
    txb.word_wrap = False
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in code.strip().split('\n'):
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(10.5)
        run.font.name = 'Consolas'
        run.font.color.rgb = GREEN

def bullet_box(sl, items, x, y, w, h, title=None, color=WHITE,
               title_color=BLUE, size=15):
    """A box with a title and bullet list."""
    box(sl, x, y, w, h, bg=RGBColor(0x24,0x24,0x37), border=GRAY)
    yo = y + 0.12
    if title:
        txt(sl, title, x+0.15, yo, w-0.3, 0.35,
            size=14, bold=True, color=title_color)
        yo += 0.38
    txb = sl.shapes.add_textbox(Inches(x+0.15), Inches(yo),
                                 Inches(w-0.3), Inches(h-(yo-y)-0.1))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = ("• " if not item.startswith("  ") else "") + item
        run.font.size = Pt(size)
        run.font.color.rgb = color

def complexity_badge(sl, text, x, y):
    """Small teal complexity label."""
    b = box(sl, x, y, 2.8, 0.32, bg=ACCENT)
    txt(sl, text, x+0.1, y+0.02, 2.6, 0.28, size=12, bold=True,
        color=DARK, align=PP_ALIGN.CENTER)

def chapter_divider(sl, num, title, desc):
    """Full-screen chapter title slide."""
    box(sl, 0, 0, 13.33, 7.5, bg=DARK)
    box(sl, 0, 2.8, 0.18, 1.9, bg=ACCENT)
    txt(sl, f"Chapter {num}", 0.5, 2.5, 12, 0.7,
        size=18, color=ACCENT, bold=True)
    txt(sl, title, 0.5, 3.1, 12, 1.1,
        size=40, bold=True, color=PURPLE)
    txt(sl, desc, 0.5, 4.3, 10, 0.8, size=16, color=WHITE)
    txt(sl, "SEng2022 — Design & Analysis of Data Structures & Algorithms",
        0.5, 6.8, 12, 0.4, size=11, color=GRAY)

def note_box(sl, text, x, y, w, color=YELLOW):
    box(sl, x, y, w, 0.45, bg=RGBColor(0x2A,0x2A,0x1E), border=color)
    txt(sl, "⚡ " + text, x+0.12, y+0.04, w-0.24, 0.37,
        size=12, color=color)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
box(sl, 0, 0, 13.33, 7.5, bg=DARK)
box(sl, 0, 0, 0.22, 7.5, bg=ACCENT)
box(sl, 0.22, 3.1, 13.11, 0.06, bg=ACCENT)
txt(sl, "Smart Campus Asset Manager", 1, 1.0, 11, 1.1,
    size=38, bold=True, color=PURPLE, align=PP_ALIGN.LEFT)
txt(sl, "University Asset Management System", 1, 2.1, 11, 0.6,
    size=22, color=WHITE, align=PP_ALIGN.LEFT)
txt(sl, "DSA Course Project  ·  SEng2022  ·  Debre Berhan University",
    1, 3.4, 11, 0.5, size=15, color=BLUE, align=PP_ALIGN.LEFT)
txt(sl, "Built With:  C++17  ·  Qt 6  ·  Oracle XE  ·  DSAEngine",
    1, 4.0, 11, 0.45, size=13, color=GRAY, align=PP_ALIGN.LEFT)
txt(sl, "Covers all 9 Chapters  ·  Chapters 1–9  ·  CLO 1–7",
    1, 4.55, 11, 0.4, size=13, color=GREEN, align=PP_ALIGN.LEFT)
for i,lab in enumerate(["Stack","Linked List","BST","Graph","Merge Sort","Greedy"]):
    box(sl, 1+i*2.0, 5.6, 1.8, 0.38, bg=ACCENT)
    txt(sl, lab, 1.05+i*2.0, 5.63, 1.7, 0.32, size=12, bold=True,
        color=DARK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 2 — TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Table of Contents")
cols = [
    [("01","Project Overview"),("02","Oracle DB & Schema"),
     ("03","Ch.1 Complexity Analysis"),("04","Ch.2 Bubble & Selection Sort"),
     ("05","Ch.3 Linked List")],
    [("06","Ch.4 Stack & Queue"),("07","Ch.5 Recursion"),
     ("08","Ch.6 Binary Search Tree"),("09","Ch.7 Graph BFS/DFS"),
     ("10","Ch.8 Merge Sort")],
    [("11","Ch.9 Greedy / D&C"),("12","Oracle Reset & Test"),
     ("13","Security Testing"),("14","Live Demo Guide"),
     ("15","Conclusion & Grade")]
]
for ci,col in enumerate(cols):
    bx = 0.4 + ci*4.3
    for ri,(num,lab) in enumerate(col):
        by = 1.3 + ri*1.1
        box(sl, bx, by, 4.0, 0.82, bg=RGBColor(0x24,0x24,0x37), border=GRAY)
        txt(sl, num, bx+0.12, by+0.08, 0.5, 0.38, size=18,
            bold=True, color=ACCENT)
        txt(sl, lab, bx+0.65, by+0.14, 3.2, 0.55, size=13, color=WHITE)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 3 — PROJECT OVERVIEW
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Project Overview",
          "Smart Campus Asset Manager — what it does and who uses it")
roles = [("ADMIN","Full access: add/delete assets, manage users"),
         ("TECHNICIAN","Maintenance + asset view"),
         ("SECURITY","Campus exit/return check"),
         ("AUDITOR","Read-only asset view"),
         ("STUDENT","My items + maintenance reports")]
for i,(role,desc) in enumerate(roles):
    bx = 0.35 + (i%3)*4.3
    by = 1.3 + (i//3)*1.6
    box(sl, bx, by, 4.0, 1.35, bg=RGBColor(0x24,0x24,0x37), border=ACCENT)
    txt(sl, role, bx+0.15, by+0.12, 3.7, 0.4, size=15, bold=True, color=PURPLE)
    txt(sl, desc, bx+0.15, by+0.55, 3.7, 0.7, size=12, color=WHITE)
note_box(sl, "All 5 roles use the same app — access is controlled by role-based policy in widget.cpp", 0.35, 6.75, 12.6)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 4 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "System Architecture",
          "How the UI, DSAEngine, and Oracle DB connect")
layers = [
    ("Qt 6 GUI Layer", "widget.cpp · assetwindow.cpp · securitycheckwindow.cpp\nmaintenancewindow.cpp · campusgraphwindow.cpp", PURPLE),
    ("DSAEngine  (dsaengine.h / dsaengine.cpp)", "Stack · Linked List · BST · Graph · Merge Sort · Greedy Sort\nAll pure C++ — no Qt dependency", BLUE),
    ("Oracle XE Database  (localhost:1522/xe)", "Tables: ASSETS · STUDENT_ITEMS · MAINTENANCE_LOG · USERS · LOST_FOUND\nStored Procs: SP_REPORT_ISSUE · SP_RESOLVE_MAINTENANCE", YELLOW),
]
for i,(title,desc,col) in enumerate(layers):
    by = 1.35 + i*1.75
    box(sl, 0.4, by, 12.5, 1.45, bg=RGBColor(0x20,0x20,0x30), border=col)
    box(sl, 0.4, by, 0.12, 1.45, bg=col)
    txt(sl, title, 0.7, by+0.12, 10, 0.42, size=16, bold=True, color=col)
    txt(sl, desc,  0.7, by+0.6,  10, 0.75, size=12, color=WHITE)
    if i < 2:
        txt(sl, "↕ calls", 6.0, by+1.5, 1.5, 0.3, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 5 — ORACLE DB SCHEMA
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Oracle Database Schema",
          "Tables used by the Smart Campus Manager — Oracle XE  localhost:1522/xe")
tables = [
    ("ASSETS", "ASSET_ID(PK) · ASSET_NAME · CATEGORY\nSERIAL_NUMBER · PURCHASE_VALUE · CONDITION\nLOCATION_ID(FK) · CREATED_BY · UPDATED_AT"),
    ("STUDENT_ITEMS", "ITEM_ID(PK) · STUDENT_ID(FK) · ITEM_NAME\nCATEGORY · BRAND · SERIAL_NUMBER · STATUS\nSTATUS: REGISTERED | OUT_OF_CAMPUS"),
    ("MAINTENANCE_LOG", "MAINTENANCE_ID(PK) · ASSET_ID(FK) · REPORTED_BY(FK)\nDESCRIPTION · PRIORITY · STATUS · REPORTED_AT\nPRIORITY: URGENT|HIGH|NORMAL|LOW"),
    ("USERS", "USER_ID(PK) · FULL_NAME · STUDENT_ID\nEMAIL · ROLE · PASSWORD_HASH"),
    ("LOST_FOUND", "LF_ID(PK) · REPORTED_BY(FK) · ITEM_TYPE\nDESCRIPTION · LOCATION · STATUS · REPORTED_AT"),
    ("LOCATIONS", "LOCATION_ID(PK) · BUILDING_NAME\nROOM_NUMBER · FLOOR"),
]
for i,( tname, tcols) in enumerate(tables):
    bx = 0.35 + (i%3)*4.3
    by = 1.3 + (i//3)*2.6
    box(sl, bx, by, 4.0, 2.35, bg=RGBColor(0x1E,0x1E,0x2E), border=BLUE)
    box(sl, bx, by, 4.0, 0.42, bg=BLUE)
    txt(sl, tname, bx+0.12, by+0.05, 3.8, 0.34, size=14, bold=True,
        color=DARK, align=PP_ALIGN.CENTER)
    txt(sl, tcols, bx+0.12, by+0.52, 3.78, 1.7, size=10.5, color=WHITE)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 6 — CHAPTER 1 DIVIDER
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
chapter_divider(sl, 1, "Complexity Analysis",
    "Big-O · Omega · Theta · Asymptotic Notation · Complexity Classes")

# ═══════════════════════════════════════════════════════════════
#  SLIDE 7 — BIG-O NOTATION
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Asymptotic Notation",
          "How we measure algorithm efficiency — independent of hardware")
notations = [
    ("O (Big-O)","Upper bound — WORST case\nO(n²) means: no worse than n² operations","f38ba8"),
    ("Ω (Omega)","Lower bound — BEST case\nΩ(n) means: at least n operations","a6e3a1"),
    ("Θ (Theta)","Tight bound — EXACT growth\nΘ(n log n) = both upper and lower","89b4fa"),
]
for i,(sym,desc,hexc) in enumerate(notations):
    bx = 0.4 + i*4.3
    col = RGBColor(int(hexc[:2],16),int(hexc[2:4],16),int(hexc[4:],16))
    box(sl, bx, 1.3, 4.0, 2.2, bg=RGBColor(0x22,0x22,0x33), border=col)
    txt(sl, sym, bx+0.15, 1.42, 3.7, 0.5, size=22, bold=True, color=col)
    txt(sl, desc, bx+0.15, 1.95, 3.7, 1.4, size=13, color=WHITE)
orders = [("O(1)","Constant","Stack push/pop"),
          ("O(log n)","Logarithmic","BST search"),
          ("O(n)","Linear","Linear search"),
          ("O(n log n)","Log-linear","Merge Sort"),
          ("O(n²)","Quadratic","Bubble Sort"),]
txt(sl, "Common Complexity Classes in Our Project:", 0.4, 3.75, 12, 0.4,
    size=14, bold=True, color=PURPLE)
for i,(big,name,where) in enumerate(orders):
    bx = 0.4+i*2.55
    box(sl, bx, 4.2, 2.38, 1.5, bg=RGBColor(0x1A,0x1A,0x2A), border=GRAY)
    txt(sl, big,   bx+0.1, 4.28, 2.2, 0.42, size=18, bold=True, color=YELLOW)
    txt(sl, name,  bx+0.1, 4.72, 2.2, 0.35, size=12, color=WHITE)
    txt(sl, where, bx+0.1, 5.07, 2.2, 0.5,  size=11, color=BLUE, italic=True)
note_box(sl, "In our code: every algorithm in dsaengine.cpp has its Big-O in the comment header above it", 0.4, 6.0, 12.5)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 8 — COMPLEXITY IN OUR PROJECT (TABLE)
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Complexity in Our Project",
          "Every DSAEngine operation — time and space complexity")
headers = ["Algorithm / Operation","Best","Average","Worst","Space","Location in Code"]
col_w   = [3.6, 1.1, 1.3, 1.3, 1.0, 4.4]
col_x   = [0.3]
for w in col_w[:-1]: col_x.append(col_x[-1]+w)
by = 1.3
box(sl, 0.3, by, 12.7, 0.42, bg=ACCENT)
for ci,h in enumerate(headers):
    txt(sl, h, col_x[ci]+0.05, by+0.05, col_w[ci]-0.1, 0.32,
        size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
rows = [
    ("Bubble Sort","O(n²)","O(n²)","O(n²)","O(1)","bubbleSortById/Value()"),
    ("Selection Sort","O(n²)","O(n²)","O(n²)","O(1)","selectionSortByName/Location()"),
    ("Merge Sort","O(n log n)","O(n log n)","O(n log n)","O(n)","mergeSortByValue/Name()"),
    ("Stack push/pop","O(1)","O(1)","O(1)*","O(n)","pushUndo() / popUndo()"),
    ("LL Enqueue","O(1)","O(1)","O(1)","O(n)","llEnqueue()"),
    ("LL Dequeue","O(1)","O(1)","O(1)","O(n)","llDequeue()"),
    ("BST Insert","O(log n)","O(log n)","O(n)**","O(n)","bstInsertHelper()"),
    ("BST Search","O(log n)","O(log n)","O(n)**","O(1)","bstSearchHelper()"),
    ("Graph BFS","O(V+E)","O(V+E)","O(V+E)","O(V)","graphBFS()"),
    ("Graph DFS","O(V+E)","O(V+E)","O(V+E)","O(V)","graphDFS()"),
]
row_colors = [RGBColor(0x1E,0x1E,0x2E), RGBColor(0x22,0x22,0x30)]
for ri,row in enumerate(rows):
    by2 = by + 0.42 + ri*0.48
    box(sl, 0.3, by2, 12.7, 0.46, bg=row_colors[ri%2])
    for ci,cell in enumerate(row):
        col = WHITE
        if ci in (1,2,3):
            if "n²" in cell: col=RED
            elif "log n" in cell or "V+E" in cell: col=GREEN
            elif "1)" in cell: col=BLUE
        txt(sl, cell, col_x[ci]+0.05, by2+0.06, col_w[ci]-0.1, 0.34,
            size=10, color=col, align=PP_ALIGN.CENTER)
txt(sl,"* Capped at MAX_STACK_SIZE=20  ** Only when all keys inserted in sorted order",
    0.3, 6.75, 12.7, 0.35, size=9, color=GRAY)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 9 — CHAPTER 2 DIVIDER
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
chapter_divider(sl, 2, "Simple Sorting & Searching",
    "Bubble Sort · Selection Sort · Linear Search  —  O(n²) algorithms")

# ═══════════════════════════════════════════════════════════════
#  SLIDE 10 — BUBBLE SORT CONCEPT
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Bubble Sort — Concept",
          "Repeatedly swap adjacent elements if out of order  ·  O(n²)")
txt(sl, "Trace  [64, 34, 25, 12, 22]  →  sorted ascending by ID:",
    0.4, 1.25, 12, 0.4, size=13, color=YELLOW)
passes = [
    ("Pass 1", "[34, 25, 12, 22, 64]", "64 bubbled to end"),
    ("Pass 2", "[25, 12, 22, 34, 64]", "34 bubbled to position"),
    ("Pass 3", "[12, 22, 25, 34, 64]", "25 in place"),
    ("Pass 4", "[12, 22, 25, 34, 64]", "Sorted!"),
]
for i,(pas,arr,note) in enumerate(passes):
    bx = 0.4 + (i%2)*6.4
    by = 1.75 + (i//2)*1.45
    box(sl, bx, by, 6.0, 1.25, bg=RGBColor(0x20,0x20,0x30), border=GRAY)
    txt(sl, pas,  bx+0.15, by+0.1,  1.2, 0.4, size=13, bold=True, color=ACCENT)
    txt(sl, arr,  bx+0.15, by+0.5,  5.5, 0.42, size=14, color=GREEN, bold=True)
    txt(sl, note, bx+0.15, by+0.88, 5.5, 0.3,  size=11, color=GRAY, italic=True)
complexity_badge(sl, "Time: O(n²)  ·  Space: O(1)  ·  In-place", 0.4, 4.85)
note_box(sl, "In our app: bubbleSortById() and bubbleSortByValue() in dsaengine.cpp — used in Asset Manager sort dropdown", 0.4, 5.5, 12.5)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 11 — BUBBLE SORT CODE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Bubble Sort — Code",
          "dsaengine.cpp  ·  bubbleSortById()  ·  exact code from our project")
code_block(sl, """\
// BUBBLE SORT — O(n²) time, O(1) space
// Compare adjacent pairs; swap if out of order.
// After each outer pass, largest element is at end.
void DSAEngine::bubbleSortById(std::vector<AssetRecord> &assets)
{
    int n = (int)assets.size();
    for (int i = 0; i < n - 1; i++) {           // n-1 passes
        for (int j = 0; j < n - i - 1; j++) {   // shrinks each pass
            if (assets[j].id > assets[j+1].id) {
                AssetRecord tmp = assets[j];     // swap
                assets[j]       = assets[j+1];
                assets[j+1]     = tmp;
            }
        }
    }
}""", 0.4, 1.25, 7.8, 4.2, title="dsaengine.cpp")
bullet_box(sl, [
    "Outer loop: runs n-1 times",
    "Inner loop: reduces by i each pass",
    "Total comparisons: n(n-1)/2  →  O(n²)",
    "No extra memory — swap in-place  →  O(1) space",
    "Used by: 'ID (Bubble Sort)' and",
    "  'Value (Bubble Sort)' in sort dropdown",
    "Located: assetwindow.cpp → sortAssets()",
], 8.4, 1.25, 4.6, 4.2, title="Key Points", size=13)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 12 — SELECTION SORT
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Selection Sort — Code",
          "dsaengine.cpp  ·  selectionSortByName()  ·  O(n²) but fewer swaps")
code_block(sl, """\
// SELECTION SORT — O(n²) time, O(1) space
// Find the minimum in the unsorted portion,
// then swap it into its correct position.
void DSAEngine::selectionSortByName(
        std::vector<AssetRecord> &assets)
{
    int n = (int)assets.size();
    for (int i = 0; i < n - 1; i++) {
        int minIdx = i;                    // assume current is min
        for (int j = i + 1; j < n; j++) {
            if (assets[j].name < assets[minIdx].name)
                minIdx = j;               // found new minimum
        }
        if (minIdx != i) {                // swap only if needed
            AssetRecord tmp   = assets[i];
            assets[i]         = assets[minIdx];
            assets[minIdx]    = tmp;
        }
    }
}""", 0.4, 1.25, 7.8, 4.5)
bullet_box(sl, [
    "Same O(n²) as Bubble Sort",
    "Key difference: at most n-1 swaps",
    "  (Bubble Sort can swap n²/2 times)",
    "Used by: 'Name (Selection Sort)'",
    "  and 'Location (Selection Sort)'",
    "  in Asset Manager sort dropdown",
    "Also: sortCombo in assetwindow.cpp",
    "  calls m_dsa.selectionSortByName()",
], 8.4, 1.25, 4.6, 4.5, title="vs Bubble Sort", size=13)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 13 — LINEAR SEARCH
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Linear Search",
          "securitycheckwindow.cpp  ·  SQL LIKE = linear scan  ·  O(n)")
code_block(sl, """\
-- Oracle SQL in SecurityCheckWindow::searchItem()
-- This is linear search — scans every row
SELECT si.ITEM_ID, u.FULL_NAME, u.STUDENT_ID,
       si.ITEM_NAME, si.CATEGORY, si.BRAND, si.STATUS
FROM   STUDENT_ITEMS si
JOIN   USERS u ON si.STUDENT_ID = u.USER_ID
WHERE  LOWER(si.SERIAL_NUMBER) = LOWER(:serial)
    OR LOWER(si.ITEM_NAME)    LIKE LOWER(:likeSerial)
    OR LOWER(si.BRAND)        LIKE LOWER(:likeSerial)
-- Oracle scans every row in STUDENT_ITEMS = O(n)
-- Our BST pre-check does the same in O(log n)""", 0.4, 1.25, 7.8, 3.5)
bullet_box(sl, [
    "SQL WHERE with LIKE = full table scan",
    "Oracle evaluates every row  →  O(n)",
    "For 1000 items: up to 1000 comparisons",
    "Our BST index cuts this to O(log n)",
    "  = ~10 comparisons for 1000 items",
    "BST result shown in blue bar",
    "  in SecurityCheckWindow UI",
], 8.4, 1.25, 4.6, 3.5, title="Linear vs BST", size=13)
note_box(sl, "The BST pre-check in our app finds serial in O(log n) BEFORE the DB query runs — shown live in Security Check window", 0.4, 5.1, 12.5)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 14 — CHAPTER 3 DIVIDER
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
chapter_divider(sl, 3, "Linked Lists",
    "Custom Singly-Linked List  ·  Dynamic memory  ·  O(1) enqueue/dequeue")

# ═══════════════════════════════════════════════════════════════
#  SLIDE 15 — LINKED LIST CONCEPT + DIAGRAM
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Linked List — Concept",
          "Each node stores data + a pointer to the next node  ·  dynamic memory")
for i,lab in enumerate(["HEAD","ADD: Laptop","EDIT: Chair","DELETE: Desk","null"]):
    bx = 0.5 + i*2.5
    col = ACCENT if i==0 else (RED if i==4 else RGBColor(0x28,0x28,0x38))
    box(sl, bx, 1.8, 2.0, 0.72, bg=col, border=GRAY)
    txt(sl, lab, bx+0.1, 1.88, 1.8, 0.52, size=13, bold=(i==0),
        color=DARK if i==0 else WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, "→", bx+2.05, 2.0, 0.35, 0.3, size=20, color=ACCENT)
bullet_box(sl, [
    "HEAD pointer: first node in list",
    "TAIL pointer: last node (for O(1) append)",
    "Each LLNode holds: data (string) + next pointer",
    "Enqueue (append at tail): O(1) — just update tail->next",
    "Dequeue (remove from head): O(1) — advance HEAD",
    "No wasted memory — grows/shrinks dynamically",
    "vs std::queue: same performance, but we own the memory",
], 0.4, 2.9, 6.2, 3.5, title="Why Linked List?", size=13)
code_block(sl, """\
struct LLNode {
    std::string data;
    LLNode     *next = nullptr;
    explicit LLNode(const std::string &d) : data(d) {}
};
// HEAD and TAIL pointers in DSAEngine:
LLNode *m_llHead = nullptr;
LLNode *m_llTail = nullptr;""", 6.8, 2.9, 6.1, 2.5)
note_box(sl, "File: dsaengine.h — struct LLNode  ·  Used as the Action Log in dsaengine.cpp", 0.4, 6.6, 12.5)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 16 — LINKED LIST CODE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Linked List — Code",
          "dsaengine.cpp  ·  llEnqueue() and llDequeue()")
code_block(sl, """\
// ENQUEUE — append at tail  O(1)
void DSAEngine::llEnqueue(const std::string &action)
{
    LLNode *node = new LLNode(action);   // allocate O(1)
    if (!m_llTail) {
        m_llHead = m_llTail = node;      // list was empty
    } else {
        m_llTail->next = node;           // link at tail O(1)
        m_llTail       = node;
    }
    ++m_llCount;
}

// DEQUEUE — remove from head  O(1)
std::string DSAEngine::llDequeue()
{
    if (!m_llHead) return "";
    std::string val = m_llHead->data;
    LLNode *old = m_llHead;
    m_llHead    = m_llHead->next;        // advance head O(1)
    if (!m_llHead) m_llTail = nullptr;
    delete old;                          // free memory
    --m_llCount;
    return val;
}""", 0.4, 1.25, 7.8, 5.2)
bullet_box(sl, [
    "Enqueue called in assetwindow.cpp:",
    "  m_dsa.enqueueAction('ADD: '+name)",
    "  m_dsa.enqueueAction('EDIT: '+name)",
    "  m_dsa.enqueueAction('DELETE: '+name)",
    "",
    "Dequeue available for audit log export",
    "",
    "llContents() shows all actions",
    "  top-to-bottom for report screen",
    "",
    "Memory freed in ~DSAEngine()",
    "  destructor — no memory leak",
], 8.4, 1.25, 4.6, 5.2, title="Where It's Called", size=12)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 17 — CHAPTER 4 DIVIDER
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
chapter_divider(sl, 4, "Stacks & Queues",
    "Stack (LIFO) for Undo/Redo  ·  Queue via custom Linked List for Action Log")

# ═══════════════════════════════════════════════════════════════
#  SLIDE 18 — STACK CONCEPT + UNDO/REDO FLOW
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Stack — Undo/Redo System",
          "LIFO (Last In First Out)  ·  O(1) push and pop")
# undo stack diagram
for i,lab in enumerate(["Add Laptop","Edit Chair","Delete Desk"]):
    col_i = [GRAY, RGBColor(0x28,0x38,0x28), ACCENT][i]
    box(sl, 0.4, 4.5-i*1.0, 3.5, 0.82, bg=col_i, border=GREEN)
    txt(sl, lab, 0.55, 4.6-i*1.0, 3.2, 0.6, size=13, color=WHITE, bold=(i==2))
txt(sl, "UNDO STACK", 0.4, 1.35, 3.5, 0.4, size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txt(sl, "← TOP (pop here)", 3.95, 2.35, 3.0, 0.4, size=11, color=ACCENT, italic=True)
txt(sl, "← push new action here", 3.95, 3.35, 3.0, 0.4, size=11, color=GRAY, italic=True)
# flow
bullet_box(sl, [
    "1. User ADDS asset  → pushUndo(ADD snapshot)",
    "2. User EDITS asset → pushUndo(EDIT snapshot)",
    "3. User clicks UNDO → popUndo() returns EDIT",
    "   → reverses the edit in Oracle DB",
    "   → pushRedo(snap) for redo",
    "4. User clicks REDO → popRedo() re-applies",
    "",
    "MAX_STACK_SIZE = 20  → bounded memory",
    "clearRedo() called on every new action",
], 4.2, 1.3, 8.8, 4.5, title="Undo/Redo Flow", size=13)
complexity_badge(sl, "push O(1)  ·  pop O(1)  ·  Space O(n) where n ≤ 20", 0.4, 6.0)
note_box(sl, "See 'Undo: N | Redo: N' counter in Asset Manager top-right — live in the app", 3.5, 6.0, 9.5)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 19 — STACK CODE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Stack — Code",
          "dsaengine.cpp  ·  pushUndo() / popUndo()  ·  uses std::stack internally")
code_block(sl, """\
// PUSH — O(1) amortized (O(n) only at MAX_STACK_SIZE cap)
void DSAEngine::pushUndo(const AssetSnapshot &snap)
{
    // Cap at 20 — drain oldest entry if full
    if ((int)m_undoStack.size() >= MAX_STACK_SIZE) {
        std::stack<AssetSnapshot> temp;
        while (m_undoStack.size() > 1) {
            temp.push(m_undoStack.top());
            m_undoStack.pop();
        }
        m_undoStack.pop(); // remove oldest
        while (!temp.empty()) {
            m_undoStack.push(temp.top());
            temp.pop();
        }
    }
    m_undoStack.push(snap);   // O(1)
}

// POP — O(1)
AssetSnapshot DSAEngine::popUndo()
{
    AssetSnapshot snap = m_undoStack.top();
    m_undoStack.pop();
    return snap;
}""", 0.4, 1.25, 7.8, 5.3)
bullet_box(sl, [
    "AssetSnapshot stores:",
    "  assetId, assetName, category",
    "  condition, serialNumber",
    "  purchaseValue, action type",
    "",
    "ActionType enum:",
    "  ADD_ASSET",
    "  DELETE_ASSET",
    "  EDIT_ASSET",
    "",
    "undoAction() in assetwindow.cpp",
    "  reverses each action type",
    "  by calling matching Oracle SQL",
], 8.4, 1.25, 4.6, 5.3, title="AssetSnapshot", size=12)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 20 — CHAPTER 5 DIVIDER
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
chapter_divider(sl, 5, "Recursion",
    "A function that calls itself  ·  Used in BST + DFS + Merge Sort")

# ═══════════════════════════════════════════════════════════════
#  SLIDE 21 — RECURSION CONCEPT
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Recursion in Our Project",
          "5 recursive functions — all in dsaengine.cpp")
uses = [
    ("bstInsertHelper()","BST insert\nGoes left/right recursively\nuntil empty node found"),
    ("bstSearchHelper()","BST search\nCompares serial, goes\nleft or right each step"),
    ("bstClearHelper()","BST memory cleanup\nPost-order traversal\ndelete left→right→node"),
    ("bstHeightHelper()","BST height\nmax(left,right)+1\nbase case = null → 0"),
    ("graphDFS()","Graph DFS\nVisit node, recurse into\neach unvisited neighbour"),
]
for i,(fname,desc) in enumerate(uses):
    bx = 0.3 + (i%3)*4.35
    by = 1.3  + (i//3)*2.2
    box(sl, bx, by, 4.15, 2.0, bg=RGBColor(0x20,0x20,0x30), border=PURPLE)
    txt(sl, fname, bx+0.12, by+0.1, 3.9, 0.42, size=13, bold=True, color=PURPLE)
    txt(sl, desc,  bx+0.12, by+0.58, 3.9, 1.25, size=12, color=WHITE)
note_box(sl, "Recursion rule: every recursive function has a BASE CASE (stops the recursion) and a RECURSIVE CASE (calls itself with smaller input)", 0.3, 6.6, 12.7)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 22 — RECURSION CODE (BST INSERT)
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Recursion — BST Insert Code",
          "dsaengine.cpp  ·  bstInsertHelper()  ·  recursive divide & search")
code_block(sl, """\
BSTNode* DSAEngine::bstInsertHelper(
        BSTNode *node, const std::string &serial, int itemId)
{
    // BASE CASE: empty spot found — create new node
    if (!node)
        return new BSTNode(serial, itemId);

    // RECURSIVE CASE: go left or right
    if (serial < node->serial)
        node->left  = bstInsertHelper(node->left,  serial, itemId);
    else if (serial > node->serial)
        node->right = bstInsertHelper(node->right, serial, itemId);
    // duplicate serial: ignore

    return node;   // return (possibly updated) node
}

// Call stack example for inserting "SN-005":
//   bstInsertHelper(root, "SN-005", 42)
//     → serial > root.serial → go right
//     → bstInsertHelper(root->right, "SN-005", 42)
//       → node is null → return new BSTNode("SN-005", 42)""",
    0.4, 1.25, 12.5, 5.3)
complexity_badge(sl, "Each call reduces search space by ~half  →  O(log n) depth", 0.4, 6.7)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 23 — CHAPTER 6 DIVIDER
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
chapter_divider(sl, 6, "Trees — Binary Search Tree",
    "BST for serial number indexing  ·  O(log n) search  ·  in-order traversal")

# ═══════════════════════════════════════════════════════════════
#  SLIDE 24 — BST DIAGRAM
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Binary Search Tree — Structure",
          "Left < Root < Right  ·  guarantees O(log n) search on average")
nodes = [("SN-050",5.9,1.5),
         ("SN-020",3.2,2.6),("SN-080",8.6,2.6),
         ("SN-010",1.5,3.7),("SN-035",4.9,3.7),("SN-065",7.2,3.7),("SN-095",10.3,3.7)]
for (lab,bx,by) in nodes:
    box(sl, bx, by, 1.7, 0.55, bg=ACCENT, border=GREEN)
    txt(sl, lab, bx+0.08, by+0.08, 1.54, 0.38, size=12, bold=True,
        color=DARK, align=PP_ALIGN.CENTER)
edges = [(5.9,1.5,3.2,2.6),(5.9,1.5,8.6,2.6),
         (3.2,2.6,1.5,3.7),(3.2,2.6,4.9,3.7),
         (8.6,2.6,7.2,3.7),(8.6,2.6,10.3,3.7)]
from pptx.util import Inches as I2
from pptx.oxml.ns import qn
import lxml.etree as etree
bullet_box(sl, [
    "Property: left subtree < node < right subtree",
    "Insert SN-040: SN-040>SN-020, SN-040>SN-035",
    "  → goes right of SN-035",
    "Search SN-035: root→left→right = 3 comparisons",
    "In-order traversal = sorted serial numbers",
    "Height = log₂(7) ≈ 3  →  max 3 comparisons",
    "Used in: SecurityCheckWindow.loadSerialsBST()",
], 0.4, 4.5, 12.5, 2.5, title="BST Rules", size=13)
note_box(sl, "Blue info bar in Security Check shows: 'BST: FOUND in O(log n) ≈ X comparisons'", 0.4, 7.0, 12.5)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 25 — BST CODE (SEARCH)
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "BST — Search & Load Code",
          "securitycheckwindow.cpp + dsaengine.cpp")
code_block(sl, """\
// ── LOAD all serials into BST on window open ──
void SecurityCheckWindow::loadSerialsBST()
{
    m_dsa.bstClear();
    QSqlQuery q(QSqlDatabase::database());
    q.exec("SELECT ITEM_ID, SERIAL_NUMBER "
           "FROM STUDENT_ITEMS WHERE SERIAL_NUMBER IS NOT NULL");
    int loaded = 0;
    while (q.next()) {
        int     id  = q.value(0).toInt();
        QString ser = q.value(1).toString().trimmed().toLower();
        if (!ser.isEmpty()) { m_dsa.bstInsert(ser.toStdString(), id); ++loaded; }
    }
    int h = m_dsa.bstHeight();
    bstLabel->setText(
        QString("BST Index: %1 serials | Height: %2 | "
                "Search cost: O(log n) ≈ %3 comparisons")
            .arg(loaded).arg(h).arg((int)ceil(log2(loaded+1))));
}

// ── SEARCH during serial lookup ──
bool bstHit = m_dsa.bstSearch(serial.toLower().toStdString(), bstItemId);
if (bstHit)
    bstLabel->setText("BST: FOUND in O(log n) ≈ " + logN + " comparisons");""",
    0.4, 1.25, 12.5, 5.3)
complexity_badge(sl, "Build BST: O(n log n)  ·  Each search: O(log n)  ·  Space: O(n)", 0.4, 6.75)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 26 — CHAPTER 7 DIVIDER
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
chapter_divider(sl, 7, "Graphs — BFS & DFS",
    "Campus Building Graph  ·  10 nodes  ·  12 edges  ·  BFS shortest path  ·  DFS traversal")

# ═══════════════════════════════════════════════════════════════
#  SLIDE 27 — GRAPH ADJACENCY LIST
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Campus Graph — Adjacency List",
          "10 DBU campus buildings as nodes  ·  12 undirected edges as paths")
adj = [
    ("Main Gate",             "Library, Admin Block"),
    ("Library",               "Main Gate, College of Computing, Cafeteria"),
    ("Admin Block",           "Main Gate, College of Computing, Lab Block"),
    ("College of Computing",  "Library, Admin Block, Lab Block"),
    ("Lab Block",             "Admin Block, College of Computing"),
    ("Cafeteria",             "Library, Sports Complex, Dormitory A"),
    ("Sports Complex",        "Cafeteria, Dormitory B"),
    ("Dormitory A",           "Cafeteria, Medical Center"),
    ("Dormitory B",           "Sports Complex, Medical Center"),
    ("Medical Center",        "Dormitory A, Dormitory B"),
]
for i,(node,edges) in enumerate(adj):
    bx = 0.35 + (i%2)*6.45
    by = 1.3  + (i//2)*1.05
    box(sl, bx, by, 6.2, 0.88, bg=RGBColor(0x1E,0x1E,0x2E), border=GRAY)
    txt(sl, node,  bx+0.12, by+0.06, 2.8, 0.38, size=12, bold=True, color=ACCENT)
    txt(sl, "→  "+edges, bx+0.12, by+0.48, 5.9, 0.32, size=10.5, color=WHITE)
note_box(sl, "graphInit() in dsaengine.cpp builds this graph at startup using graphAddEdge() — 12 calls for 12 edges", 0.35, 6.65, 12.6)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 28 — BFS CODE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Graph — BFS Shortest Path Code",
          "dsaengine.cpp  ·  graphBFS()  ·  uses std::queue + parent map")
code_block(sl, """\
std::vector<std::string> DSAEngine::graphBFS(
        const std::string &start, const std::string &end) const
{
    std::unordered_map<std::string,std::string> parent;
    std::queue<std::string>                     bfsQ;
    std::unordered_set<std::string>             visited;

    bfsQ.push(start);
    visited.insert(start);
    parent[start] = "";

    while (!bfsQ.empty()) {
        std::string node = bfsQ.front();  // FIFO — process level by level
        bfsQ.pop();

        if (node == end) {
            // Reconstruct path: walk parent pointers back to start
            std::vector<std::string> path;
            for (std::string cur=end; cur!=""; cur=parent[cur])
                path.push_back(cur);
            std::reverse(path.begin(), path.end());
            return path;
        }
        for (const std::string &nb : m_graph.at(node))
            if (!visited.count(nb)) {
                visited.insert(nb); parent[nb]=node; bfsQ.push(nb);
            }
    }
    return {};  // no path
}""", 0.4, 1.25, 8.0, 5.5)
bullet_box(sl, [
    "Uses a queue (FIFO) — links Ch.3/4",
    "parent[] map records the path",
    "Guarantees SHORTEST path in",
    "  unweighted graph",
    "Example result:",
    "  Main Gate → Library →",
    "  College of Computing",
    "  = 2 hops (shortest possible)",
    "",
    "Complexity: O(V+E) = O(10+12)",
    "V = buildings, E = paths",
], 8.5, 1.25, 4.5, 5.5, title="BFS Key Points", size=12)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 29 — DFS CODE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Graph — DFS Traversal Code",
          "dsaengine.cpp  ·  graphDFS()  ·  recursive  ·  links to Ch.5 Recursion")
code_block(sl, """\
// DFS — visit node, recurse into each unvisited neighbour
// Uses recursion (implicit call stack = Ch.5 connection)
std::vector<std::string> DSAEngine::graphDFS(
        const std::string &start) const
{
    std::vector<std::string>        visited;
    std::unordered_set<std::string> seen;

    // Lambda captures visited/seen by reference
    std::function<void(const std::string&)> dfs =
        [&](const std::string &node) {
            seen.insert(node);
            visited.push_back(node);           // record visit order
            auto it = m_graph.find(node);
            if (it == m_graph.end()) return;
            for (const std::string &nb : it->second)
                if (!seen.count(nb))
                    dfs(nb);                   // RECURSIVE CALL
        };

    if (m_graph.count(start))
        dfs(start);
    return visited;   // visit order from start
}""", 0.4, 1.25, 8.0, 5.5)
bullet_box(sl, [
    "Recursive lambda = Ch.5 recursion",
    "seen set prevents revisiting",
    "visited vector = output order",
    "",
    "BFS vs DFS:",
    "  BFS: level by level (queue)",
    "  DFS: depth first (recursion)",
    "  BFS: guaranteed shortest path",
    "  DFS: full connectivity check",
    "",
    "O(V+E) — same as BFS",
    "Used in: Campus Graph window",
    "  '▶ Run DFS' button",
], 8.5, 1.25, 4.5, 5.5, title="DFS vs BFS", size=12)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 30 — CHAPTER 8 DIVIDER
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
chapter_divider(sl, 8, "Advanced Sorting",
    "Merge Sort (Divide & Conquer)  ·  O(n log n)  ·  vs Bubble O(n²)")

# ═══════════════════════════════════════════════════════════════
#  SLIDE 31 — MERGE SORT CONCEPT
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Merge Sort — Divide & Conquer",
          "T(n) = 2T(n/2) + O(n)  →  O(n log n)  ·  Ch.8 + Ch.9 combined")
# tree diagram text representation
trace = [
    (5.7,  1.3,  2.0, "[38,27,43,3]",     PURPLE),
    (2.8,  2.1,  2.0, "[38,27]",          BLUE),
    (8.6,  2.1,  2.0, "[43,3]",           BLUE),
    (1.3,  2.9,  1.5, "[38]",             GREEN),
    (4.3,  2.9,  1.5, "[27]",             GREEN),
    (7.1,  2.9,  1.5, "[43]",             GREEN),
    (9.9,  2.9,  1.5, "[3]",              GREEN),
    (2.8,  3.75, 2.0, "[27,38]",          YELLOW),
    (8.6,  3.75, 2.0, "[3,43]",           YELLOW),
    (5.7,  4.55, 2.0, "[3,27,38,43]",     ACCENT),
]
for (bx,by,bw,lab,col) in trace:
    box(sl, bx, by, bw, 0.52, bg=RGBColor(0x20,0x20,0x30), border=col)
    txt(sl, lab, bx+0.08, by+0.06, bw-0.16, 0.4, size=12,
        bold=True, color=col, align=PP_ALIGN.CENTER)
txt(sl, "DIVIDE →", 0.3, 2.9, 1.3, 0.5, size=11, color=GRAY, italic=True)
txt(sl, "← MERGE", 0.3, 3.75,1.3, 0.5, size=11, color=GRAY, italic=True)
bullet_box(sl, [
    "DIVIDE: split in half recursively until size=1",
    "CONQUER: merge two sorted halves",
    "Recurrence: T(n) = 2T(n/2) + O(n)",
    "  → O(n log n) by Master Theorem",
    "For 500 assets:",
    "  Bubble: 500×500 = 250,000 ops",
    "  Merge:  500×9 ≈ 4,500 ops  (55× faster!)",
], 0.35, 5.25, 12.6, 2.0, title="Why Merge Sort?", size=13)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 32 — MERGE SORT CODE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Merge Sort — Code",
          "dsaengine.cpp  ·  mergeSortByValue()  ·  recursive divide & conquer")
code_block(sl, """\
// Entry point
void DSAEngine::mergeSortByValue(std::vector<AssetRecord> &assets)
{
    if (assets.size() > 1)
        mergeSortHelperByValue(assets, 0, (int)assets.size()-1);
}

// DIVIDE step — recursive
void DSAEngine::mergeSortHelperByValue(
        std::vector<AssetRecord> &a, int l, int r)
{
    if (l >= r) return;                    // BASE CASE: size 1
    int m = l + (r - l) / 2;              // find midpoint
    mergeSortHelperByValue(a, l, m);       // sort left half
    mergeSortHelperByValue(a, m+1, r);     // sort right half
    mergeByValue(a, l, m, r);             // CONQUER: merge
}

// CONQUER step — merge two sorted halves
void DSAEngine::mergeByValue(
        std::vector<AssetRecord> &a, int l, int m, int r)
{
    std::vector<AssetRecord> left (a.begin()+l,   a.begin()+m+1);
    std::vector<AssetRecord> right(a.begin()+m+1, a.begin()+r+1);
    int i=0, j=0, k=l;
    while (i<(int)left.size() && j<(int)right.size())
        a[k++] = (left[i].value <= right[j].value)
                 ? left[i++] : right[j++];
    while (i<(int)left.size())  a[k++]=left[i++];
    while (j<(int)right.size()) a[k++]=right[j++];
}""", 0.4, 1.25, 8.0, 5.7)
bullet_box(sl, [
    "UI: select 'Value (Merge Sort - D&C)'",
    "  in sort dropdown → calls:",
    "  m_dsa.mergeSortByValue(records)",
    "  in assetwindow.cpp sortAssets()",
    "",
    "Also available:",
    "  mergeSortByName() — same logic",
    "  for alphabetical sort",
    "",
    "Time: O(n log n) all cases",
    "Space: O(n) auxiliary buffer",
    "  (left/right temp vectors)",
    "",
    "Compare on same data with",
    "  Bubble Sort to show difference",
], 8.5, 1.25, 4.5, 5.7, title="Usage", size=12)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 33 — CHAPTER 9 DIVIDER
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
chapter_divider(sl, 9, "Algorithmic Strategies",
    "Greedy Scheduling  ·  Divide & Conquer  ·  Brute Force comparison")

# ═══════════════════════════════════════════════════════════════
#  SLIDE 34 — GREEDY CONCEPT
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Greedy Algorithm — Maintenance Scheduler",
          "Always pick the locally optimal choice → globally optimal schedule")
scores = [("URGENT + PENDING","40+20 = 60","#1 — Fix NOW"),
          ("URGENT + IN_PROGRESS","40+10 = 50","#2 — Already started"),
          ("HIGH + PENDING","30+20 = 50","#3 — Next in queue"),
          ("HIGH + IN_PROGRESS","30+10 = 40","#4"),
          ("NORMAL + PENDING","20+20 = 40","#5"),
          ("LOW + RESOLVED","10+0  = 10","#Last — Done")]
txt(sl, "Greedy Score = Priority Weight + Status Weight:", 0.4, 1.3, 12.5, 0.4,
    size=14, bold=True, color=YELLOW)
for i,(label,score,rank) in enumerate(scores):
    bx = 0.4 + (i%3)*4.3
    by = 1.8 + (i//3)*1.55
    col = [RED,RED,YELLOW,YELLOW,BLUE,GRAY][i]
    box(sl, bx, by, 4.1, 1.35, bg=RGBColor(0x20,0x20,0x30), border=col)
    txt(sl, label, bx+0.12, by+0.1,  3.86, 0.4,  size=13, color=col, bold=True)
    txt(sl, score, bx+0.12, by+0.55, 3.86, 0.35, size=14, color=WHITE, bold=True)
    txt(sl, rank,  bx+0.12, by+0.9,  3.86, 0.32, size=11, color=ACCENT, italic=True)
note_box(sl, "Greedy choice property: picking highest-score task first minimizes maximum waiting time for critical assets — provably optimal for this scheduling problem", 0.4, 6.6, 12.6)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 35 — GREEDY CODE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Greedy — Code",
          "maintenancewindow.cpp  ·  greedyPrioritize()  ·  ⚡ Auto-Prioritize button")
code_block(sl, """\
void MaintenanceWindow::greedyPrioritize()
{
    struct Row { int score; QStringList cols; QVector<QColor> fgColors; };
    QVector<Row> rows;

    for (int r = 0; r < table->rowCount(); r++) {
        Row row; row.score = 0;
        for (int c = 0; c < table->columnCount(); c++) {
            auto *it = table->item(r,c);
            row.cols     << (it ? it->text()           : "");
            row.fgColors << (it ? it->foreground().color() : QColor("#cdd6f4"));
        }
        // GREEDY KEY: assign urgency score
        const QString &pri = row.cols[4];
        if      (pri=="URGENT") row.score += 40;
        else if (pri=="HIGH")   row.score += 30;
        else if (pri=="NORMAL") row.score += 20;
        else                    row.score += 10;   // LOW
        // Status modifier
        const QString &st = row.cols[5];
        if      (st=="PENDING")     row.score += 20;
        else if (st=="IN_PROGRESS") row.score += 10;
        rows.append(row);
    }
    // GREEDY SELECTION: sort descending by score — O(n log n)
    std::sort(rows.begin(), rows.end(),
              [](const Row &a, const Row &b){ return a.score > b.score; });
    // Repopulate table in greedy order
    table->setRowCount(0);
    for (const Row &row : rows) { /* insert rows back */ }
}""", 0.4, 1.25, 12.5, 5.7)
complexity_badge(sl, "Time: O(n log n) for std::sort  ·  Space: O(n) for row buffer", 0.4, 7.15)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 36 — ORACLE RESET & TEST — DIVIDER
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
chapter_divider(sl, "—", "Oracle DB — Reset & Test",
    "How to connect, reset tables, test stored procedures, and verify data")

# ═══════════════════════════════════════════════════════════════
#  SLIDE 37 — ORACLE CONNECTION
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Oracle XE — Connection Details",
          "How the app connects  ·  widget.cpp")
code_block(sl, """\
-- Connection string in widget.cpp:
db.setDatabaseName(
    "Driver={Oracle in XE};"
    "DBQ=127.0.0.1:1522/xe;"
);
db.setUserName("campus_admin");
db.setPassword("campus123");

-- Connect via SQL*Plus (command line):
sqlplus campus_admin/campus123@//127.0.0.1:1522/xe

-- Or via Oracle SQL Developer:
Host:     127.0.0.1
Port:     1522
SID/Svc:  xe
Username: campus_admin
Password: campus123

-- Check connection is alive:
SELECT 'Connected as: ' || USER FROM DUAL;""", 0.4, 1.3, 8.0, 5.2)
bullet_box(sl, [
    "Oracle XE must be running before",
    "  launching the app",
    "",
    "Start Oracle service (Windows):",
    "  net start OracleServiceXE",
    "  net start OracleXETNSListener",
    "",
    "Check service status:",
    "  sc query OracleServiceXE",
    "",
    "If port 1522 is busy:",
    "  netstat -ano | findstr 1522",
], 8.5, 1.3, 4.6, 5.2, title="Oracle Startup", size=12)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 38 — ORACLE RESET COMMANDS
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Oracle — Reset / Clear Data",
          "Use these in SQL*Plus or SQL Developer to reset to clean state")
code_block(sl, """\
-- ① Disable foreign keys temporarily
ALTER TABLE MAINTENANCE_LOG DISABLE CONSTRAINT FK_MAINT_ASSET;
ALTER TABLE STUDENT_ITEMS   DISABLE CONSTRAINT FK_ITEMS_USER;
ALTER TABLE LOST_FOUND      DISABLE CONSTRAINT FK_LF_USER;

-- ② Clear child tables first (order matters!)
DELETE FROM MAINTENANCE_LOG;
DELETE FROM LOST_FOUND;
DELETE FROM STUDENT_ITEMS;

-- ③ Clear parent tables
DELETE FROM ASSETS;
DELETE FROM USERS WHERE ROLE != 'ADMIN';  -- keep admin

-- ④ Reset sequences (auto-increment IDs back to 1)
DROP SEQUENCE SEQ_ASSETS;
CREATE SEQUENCE SEQ_ASSETS START WITH 1 INCREMENT BY 1;

DROP SEQUENCE SEQ_MAINTENANCE;
CREATE SEQUENCE SEQ_MAINTENANCE START WITH 1 INCREMENT BY 1;

-- ⑤ Re-enable foreign keys
ALTER TABLE MAINTENANCE_LOG ENABLE CONSTRAINT FK_MAINT_ASSET;
ALTER TABLE STUDENT_ITEMS   ENABLE CONSTRAINT FK_ITEMS_USER;

-- ⑥ Commit everything
COMMIT;

-- ⑦ Verify counts (should all be 0 or 1 for admin)
SELECT COUNT(*) FROM ASSETS;
SELECT COUNT(*) FROM MAINTENANCE_LOG;
SELECT COUNT(*) FROM STUDENT_ITEMS;""", 0.4, 1.25, 12.5, 5.9)
note_box(sl, "Always COMMIT after deletes — Oracle does NOT auto-commit unlike MySQL", 0.4, 7.2, 12.5)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 39 — INSERT TEST DATA
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Oracle — Insert Test Data",
          "Seed the database for a full demo")
code_block(sl, """\
-- Insert test users
INSERT INTO USERS(USER_ID,FULL_NAME,STUDENT_ID,EMAIL,ROLE,PASSWORD_HASH)
VALUES(1,'Admin User','ADM001','admin@dbu.edu.et','ADMIN','admin123');
INSERT INTO USERS(USER_ID,FULL_NAME,STUDENT_ID,EMAIL,ROLE,PASSWORD_HASH)
VALUES(2,'Bereket Tadesse','DBU/1234/15','berek@dbu.edu.et','STUDENT','pass123');
INSERT INTO USERS(USER_ID,FULL_NAME,STUDENT_ID,EMAIL,ROLE,PASSWORD_HASH)
VALUES(3,'Security Guard','SEC001','sec@dbu.edu.et','SECURITY','sec123');

-- Insert test assets
INSERT INTO ASSETS(ASSET_ID,ASSET_NAME,CATEGORY,SERIAL_NUMBER,PURCHASE_VALUE,CONDITION,CREATED_BY)
VALUES(SEQ_ASSETS.NEXTVAL,'Dell Laptop 15','IT_EQUIPMENT','DL-2024-001',45000,'GOOD',1);
INSERT INTO ASSETS(ASSET_ID,ASSET_NAME,CATEGORY,SERIAL_NUMBER,PURCHASE_VALUE,CONDITION,CREATED_BY)
VALUES(SEQ_ASSETS.NEXTVAL,'Projector Epson','IT_EQUIPMENT','EP-2023-042',28000,'GOOD',1);
INSERT INTO ASSETS(ASSET_ID,ASSET_NAME,CATEGORY,SERIAL_NUMBER,PURCHASE_VALUE,CONDITION,CREATED_BY)
VALUES(SEQ_ASSETS.NEXTVAL,'Lab Microscope','LAB_EQUIPMENT','MIC-001',120000,'UNDER_MAINTENANCE',1);

-- Insert student items (for BST + Security demo)
INSERT INTO STUDENT_ITEMS(ITEM_ID,STUDENT_ID,ITEM_NAME,CATEGORY,BRAND,SERIAL_NUMBER,STATUS)
VALUES(1,2,'MacBook Pro','LAPTOP','Apple','MB-PRO-2024-XYZ','REGISTERED');
INSERT INTO STUDENT_ITEMS(ITEM_ID,STUDENT_ID,ITEM_NAME,CATEGORY,BRAND,SERIAL_NUMBER,STATUS)
VALUES(2,2,'Samsung Phone','PHONE','Samsung','SM-G998B-123','REGISTERED');

-- Insert maintenance log (for Greedy demo)
BEGIN SP_REPORT_ISSUE(1,1,'BROKEN','Screen cracked','URGENT'); END;
BEGIN SP_REPORT_ISSUE(1,2,'MISSING','Power cable lost','NORMAL'); END;

COMMIT;""", 0.4, 1.25, 12.5, 6.1)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 40 — TEST STORED PROCEDURES
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Oracle — Test Stored Procedures",
          "Verify SP_REPORT_ISSUE and SP_RESOLVE_MAINTENANCE work correctly")
code_block(sl, """\
-- ① Test SP_REPORT_ISSUE
BEGIN
    SP_REPORT_ISSUE(
        :asset_id   => 1,
        :reported_by => 1,
        :issue_type  => 'BROKEN',
        :description => 'Keyboard not working',
        :priority    => 'HIGH'
    );
END;
/

-- Verify it was inserted:
SELECT MAINTENANCE_ID, ASSET_ID, PRIORITY, STATUS, REPORTED_AT
FROM   MAINTENANCE_LOG
ORDER  BY REPORTED_AT DESC
FETCH FIRST 5 ROWS ONLY;

-- ② Test SP_RESOLVE_MAINTENANCE
BEGIN
    SP_RESOLVE_MAINTENANCE(
        :maintenance_id => 1,
        :resolved_by    => 1,
        :resolution_note => 'Replaced keyboard unit'
    );
END;
/

-- Verify status changed to RESOLVED:
SELECT STATUS, RESOLVED_AT
FROM   MAINTENANCE_LOG
WHERE  MAINTENANCE_ID = 1;

-- ③ Test STUDENT_ITEMS status update (Security Check)
UPDATE STUDENT_ITEMS SET STATUS='OUT_OF_CAMPUS' WHERE ITEM_ID=1;
COMMIT;
-- Then in the app: search 'MB-PRO-2024-XYZ' → should show OUT_OF_CAMPUS
-- Click 'Mark Returned' → run again → should show REGISTERED""",
    0.4, 1.25, 12.5, 6.1)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 41 — VERIFY QUERIES
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Oracle — Verification Queries",
          "Run these to confirm the app and DB are in sync")
code_block(sl, """\
-- ① Count all tables
SELECT 'ASSETS'          AS tbl, COUNT(*) AS cnt FROM ASSETS     UNION ALL
SELECT 'USERS',                  COUNT(*)         FROM USERS      UNION ALL
SELECT 'STUDENT_ITEMS',          COUNT(*)         FROM STUDENT_ITEMS UNION ALL
SELECT 'MAINTENANCE_LOG',        COUNT(*)         FROM MAINTENANCE_LOG UNION ALL
SELECT 'LOST_FOUND',             COUNT(*)         FROM LOST_FOUND;

-- ② Check asset conditions distribution
SELECT CONDITION, COUNT(*) AS total
FROM   ASSETS
GROUP  BY CONDITION
ORDER  BY total DESC;

-- ③ Check maintenance priority distribution (matches Greedy scores)
SELECT PRIORITY, STATUS, COUNT(*) AS total
FROM   MAINTENANCE_LOG
GROUP  BY PRIORITY, STATUS
ORDER  BY PRIORITY, STATUS;

-- ④ Check student items out of campus
SELECT si.SERIAL_NUMBER, u.FULL_NAME, si.STATUS
FROM   STUDENT_ITEMS si
JOIN   USERS u ON si.STUDENT_ID = u.USER_ID
WHERE  si.STATUS = 'OUT_OF_CAMPUS';

-- ⑤ Confirm sequences are working
SELECT SEQ_ASSETS.NEXTVAL FROM DUAL;
-- If this returns a number, sequences are active""",
    0.4, 1.25, 12.5, 5.8)
note_box(sl, "Run these queries BEFORE your demo to confirm database is in the expected state", 0.4, 7.1, 12.5)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 42 — SECURITY TESTING DIVIDER
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
chapter_divider(sl, "—", "Security Testing",
    "Test the Security Check window  ·  BST verification  ·  role-based access")

# ═══════════════════════════════════════════════════════════════
#  SLIDE 43 — SECURITY CHECK TESTING
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Security Check — Testing Guide",
          "securitycheckwindow.cpp  ·  BST + Oracle integration test")
tests = [
    ("Test 1: Registered item","Search: MB-PRO-2024-XYZ\nExpect: green label, Allow Exit enabled\nBST bar: FOUND in O(log n)"),
    ("Test 2: Out-of-campus item","First run: UPDATE STUDENT_ITEMS\nSET STATUS='OUT_OF_CAMPUS'\nWHERE ITEM_ID=1; COMMIT;\nSearch again → orange label, Mark Returned"),
    ("Test 3: Unknown serial","Search: FAKE-SN-9999\nExpect: red label 'No registered item'\nBST bar: NOT in index"),
    ("Test 4: Return item","Click 'Mark Returned'\nCheck DB: SELECT STATUS FROM\nSTUDENT_ITEMS WHERE ITEM_ID=1\nExpect: REGISTERED"),
    ("Test 5: BST accuracy","Count rows in STUDENT_ITEMS\nCompare with BST size shown\nin blue bar — must match"),
    ("Test 6: Role access","Login as STUDENT role\nSecurity Check button must\nbe HIDDEN (not in sidebar)"),
]
for i,(title,steps) in enumerate(tests):
    bx = 0.35 + (i%3)*4.3
    by = 1.3  + (i//3)*2.55
    box(sl, bx, by, 4.1, 2.35, bg=RGBColor(0x1E,0x1E,0x2E), border=BLUE)
    box(sl, bx, by, 4.1, 0.42, bg=BLUE)
    txt(sl, title, bx+0.1, by+0.06, 3.9, 0.3, size=12, bold=True,
        color=DARK, align=PP_ALIGN.CENTER)
    txt(sl, steps, bx+0.1, by+0.52, 3.9, 1.72, size=10.5, color=WHITE)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 44 — ROLE-BASED ACCESS TEST
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Role-Based Access Testing",
          "widget.cpp  ·  applyRoleAccess()  ·  RolePolicy struct")
code_block(sl, """\
-- Create test accounts for each role
INSERT INTO USERS VALUES(10,'Test Admin',   'A01','a@dbu.et','ADMIN',    'pass');
INSERT INTO USERS VALUES(11,'Test Tech',    'T01','t@dbu.et','TECHNICIAN','pass');
INSERT INTO USERS VALUES(12,'Test Security','S01','s@dbu.et','SECURITY',  'pass');
INSERT INTO USERS VALUES(13,'Test Auditor', 'AU1','u@dbu.et','AUDITOR',   'pass');
INSERT INTO USERS VALUES(14,'Test Student', 'ST1','st@dbu.et','STUDENT',  'pass');
COMMIT;

/* Expected sidebar buttons per role:
   ADMIN:      Home Assets Maintenance LostFound MyItems Users Security Graph ✓ all
   TECHNICIAN: Home Assets Maintenance LostFound MyItems             Graph
   SECURITY:   Home                   LostFound MyItems Security    Graph
   AUDITOR:    Home Assets Maintenance LostFound MyItems             Graph
   STUDENT:    Home             Maintenance LostFound MyItems        Graph      */""",
    0.4, 1.25, 8.0, 3.6)
bullet_box(sl, [
    "How to test:",
    "1. Login with each test account",
    "2. Check which sidebar buttons",
    "   are visible vs hidden",
    "3. Try navigating directly to",
    "   a blocked screen via code",
    "   → should show 'Access Denied'",
    "",
    "Code location:",
    "  widget.cpp line ~103",
    "  applyRoleAccess()",
    "  RolePolicy::forRole()",
], 8.5, 1.25, 4.6, 3.6, title="Test Steps", size=12)
note_box(sl, "showAccessDenied() in widget.cpp displays a lock screen when role lacks permission — test by calling btnAssets with SECURITY role", 0.4, 5.15, 12.5)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 45 — LIVE DEMO GUIDE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
chapter_divider(sl, "—", "Live Demo Guide",
    "Exact click-by-click sequence for a 10-minute presentation demo")

# ═══════════════════════════════════════════════════════════════
#  SLIDE 46 — DEMO SEQUENCE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Demo Click Sequence",
          "Follow this order — each step highlights a different DSA concept")
steps = [
    ("1","Login as ADMIN","Shows role-based auth → DB query in logindialog.cpp","89b4fa"),
    ("2","Asset Manager → Add asset","pushUndo() called → Undo counter increments","a6e3a1"),
    ("3","Add 2 more assets","Undo: 3 | Redo: 0 in top bar","a6e3a1"),
    ("4","Click Undo","popUndo() → Oracle DELETE → Redo: 1","f9e2af"),
    ("5","Sort dropdown → Bubble Sort ID","bubbleSortById() — O(n²) running","cba6f7"),
    ("6","Sort → Merge Sort Value","mergeSortByValue() — O(n log n) — same result, faster algo","cba6f7"),
    ("7","Maintenance → ⚡ Auto-Prioritize","greedyPrioritize() — URGENT PENDING moves to top","f38ba8"),
    ("8","Sidebar → Security Check","loadSerialsBST() runs — blue bar shows N serials loaded","89b4fa"),
    ("9","Type a registered serial → Search","BST: FOUND in O(log n) ≈ X comparisons","a6e3a1"),
    ("10","Sidebar → Campus Graph","graphInit() — adjacency list displayed","cba6f7"),
    ("11","BFS: Main Gate → Medical Center","Shows shortest path with hop count","f9e2af"),
    ("12","DFS: College of Computing","Shows all reachable buildings in DFS order","a6e3a1"),
]
for i,(num,action,explain,hexc) in enumerate(steps):
    bx = 0.35 + (i%2)*6.45
    by = 1.25 + (i//2)*0.95
    col = RGBColor(int(hexc[:2],16),int(hexc[2:4],16),int(hexc[4:],16))
    box(sl, bx, by, 6.2, 0.82, bg=RGBColor(0x1E,0x1E,0x2E), border=col)
    txt(sl, num,    bx+0.1,  by+0.08, 0.45, 0.65, size=16, bold=True, color=col)
    txt(sl, action, bx+0.6,  by+0.06, 3.4,  0.35, size=12, bold=True, color=WHITE)
    txt(sl, explain,bx+0.6,  by+0.43, 5.45, 0.32, size=10, color=GRAY, italic=True)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 47 — EXPECTED QUESTIONS & ANSWERS
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Expected Questions & Answers",
          "The 6 most likely questions from your examiner")
qas = [
    ("Why linked list instead of std::queue?",
     "We own the memory. std::queue wraps std::deque which pre-allocates chunks.\nOur LLNode allocates exactly what it needs. Also demonstrates Ch.3 directly."),
    ("BST worst case is O(n) — when?",
     "Only when all keys are inserted in sorted order (degenerate = linear chain).\nFor random serial numbers like SN-2024-XYZ this never happens in practice."),
    ("Why not hash table for serial search?",
     "Hash table gives O(1) average but loses ordering. BST in-order traversal\ngives sorted serials for free — useful for audit reports."),
    ("Why Merge Sort and not Quick Sort?",
     "Merge Sort is O(n log n) guaranteed worst case. Quick Sort degrades to O(n²)\nwith bad pivot. We chose stability over average-case performance."),
    ("How does Greedy guarantee optimal?",
     "Greedy choice property holds: URGENT tasks have higher value-per-time-unit.\nEarlier service = less damage = lower total cost. Proven optimal for this scheduling."),
    ("DFS vs BFS — when to use which?",
     "BFS = shortest path (uses queue). DFS = check connectivity, detect cycles.\nFor campus navigation (shortest route) we use BFS. DFS shows full coverage."),
]
for i,(q,a) in enumerate(qas):
    bx = 0.35 + (i%2)*6.45
    by = 1.25 + (i//2)*1.85
    box(sl, bx, by, 6.2, 1.7, bg=RGBColor(0x1E,0x1E,0x2E), border=PURPLE)
    txt(sl, "Q: "+q, bx+0.12, by+0.08, 5.96, 0.48, size=11.5, bold=True, color=PURPLE)
    txt(sl, "A: "+a, bx+0.12, by+0.6,  5.96, 0.98, size=11,   color=WHITE)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 48 — DSA COVERAGE SUMMARY
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Complete DSA Coverage — All 9 Chapters",
          "Every chapter of SEng2022 has a working implementation in the project")
rows2 = [
    ("Ch.1","Complexity Analysis","Big-O headers on every function in dsaengine.cpp","CLO2","✓"),
    ("Ch.2","Simple Sorting/Search","bubbleSort×2, selectionSort×2, linear search","CLO3","✓"),
    ("Ch.3","Linked Lists","Custom singly-LL (LLNode) as action log queue","CLO1,4","✓"),
    ("Ch.4","Stacks & Queues","std::stack for undo/redo, LL for queue","CLO1,4","✓"),
    ("Ch.5","Recursion","BST insert/search/clear/height + DFS","CLO4,5","✓"),
    ("Ch.6","Trees / BST","BST serial index in SecurityCheckWindow","CLO1,4","✓"),
    ("Ch.7","Graphs","Adjacency list + BFS + DFS, CampusGraphWindow","CLO1,4","✓"),
    ("Ch.8","Advanced Sorting","Merge Sort O(n log n), compare with bubble","CLO3","✓"),
    ("Ch.9","Algo Strategies","Greedy scheduling + Divide & Conquer (Merge Sort)","CLO5","✓"),
]
hdr = ["Chapter","Topic","Implementation in Project","CLOs","Done"]
hw  = [0.8,1.9,5.6,1.2,0.7]
hx  = [0.3]
for w in hw[:-1]: hx.append(hx[-1]+w)
by0 = 1.3
box(sl, 0.3, by0, 12.7, 0.42, bg=ACCENT)
for ci,h in enumerate(hdr):
    txt(sl, h, hx[ci]+0.05, by0+0.05, hw[ci]-0.1, 0.32,
        size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
for ri,row in enumerate(rows2):
    by2 = by0+0.42+ri*0.56
    box(sl, 0.3, by2, 12.7, 0.54,
        bg=[RGBColor(0x1E,0x1E,0x2E),RGBColor(0x22,0x22,0x30)][ri%2])
    for ci,cell in enumerate(row):
        col = GREEN if ci==4 else (YELLOW if ci==3 else WHITE)
        txt(sl, cell, hx[ci]+0.05, by2+0.08, hw[ci]-0.1, 0.38,
            size=10.5, color=col, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 49 — GRADE ESTIMATE
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
title_bar(sl, "Grade Estimate — Project Component",
          "Based on SEng2022 assessment criteria and CLO coverage")
clos = [
    ("CLO2","Complexity Analysis","Big-O on every algorithm","20/20",GREEN),
    ("CLO3","Sorting & Searching","4 sort algorithms + linear + BST","19/20",GREEN),
    ("CLO4","Apply Appropriate DS","Stack/LL/BST/Graph all motivated","19/20",GREEN),
    ("CLO5","Algorithm Strategies","Greedy + Divide & Conquer","18/20",YELLOW),
    ("CLO6","DS&A Foundation","9 chapters, real-world Qt/Oracle app","20/20",GREEN),
    ("CLO7","Team / Demo","Multi-module, live demo ready","19/20",GREEN),
]
for i,(clo,name,detail,score,col) in enumerate(clos):
    by = 1.3 + i*0.84
    box(sl, 0.35, by, 12.6, 0.76, bg=RGBColor(0x1E,0x1E,0x2E), border=col)
    txt(sl, clo,   0.5,  by+0.12, 0.9,  0.5,  size=15, bold=True, color=col)
    txt(sl, name,  1.5,  by+0.06, 3.0,  0.34, size=13, bold=True, color=WHITE)
    txt(sl, detail,1.5,  by+0.42, 7.5,  0.28, size=11, color=GRAY, italic=True)
    txt(sl, score, 11.5, by+0.1,  1.3,  0.5,  size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
box(sl, 0.35, 6.38, 12.6, 0.6, bg=ACCENT)
txt(sl, "Estimated Project Score:  115 / 120  ≈  95 %  →  A+",
    0.5, 6.44, 12.4, 0.48, size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
#  SLIDE 50 — CONCLUSION
# ═══════════════════════════════════════════════════════════════
sl = add_slide()
box(sl, 0, 0, 13.33, 7.5, bg=DARK)
box(sl, 0, 0, 0.22, 7.5, bg=ACCENT)
box(sl, 0.22, 3.6, 13.11, 0.06, bg=ACCENT)
txt(sl, "Thank You", 1, 1.2, 11, 1.0,
    size=46, bold=True, color=PURPLE, align=PP_ALIGN.LEFT)
txt(sl, "Smart Campus Asset Manager — SEng2022 Course Project",
    1, 2.25, 11, 0.5, size=18, color=WHITE, align=PP_ALIGN.LEFT)
items = [
    ("9 Chapters", "Complete Ch.1–Ch.9 coverage"),
    ("5 Data Structures", "Stack · Linked List · BST · Graph · Array"),
    ("4 Algorithms", "Bubble · Selection · Merge Sort · Greedy"),
    ("Oracle XE", "Full CRUD with stored procedures"),
    ("Qt 6 GUI", "5 role-based windows + Campus Graph"),
]
for i,(bold_t,rest) in enumerate(items):
    bx = 1 + (i%3)*3.9
    by = 3.8 + (i//3)*1.3
    box(sl, bx, by, 3.65, 1.1, bg=RGBColor(0x20,0x20,0x30), border=ACCENT)
    txt(sl, bold_t, bx+0.12, by+0.1, 3.4, 0.38, size=14, bold=True, color=ACCENT)
    txt(sl, rest,   bx+0.12, by+0.5, 3.4, 0.48, size=11, color=WHITE)
txt(sl, "Questions welcome  ·  bekelemengesha@dbu.edu.et",
    1, 6.95, 11, 0.35, size=12, color=GRAY, align=PP_ALIGN.LEFT)

# ── Save ─────────────────────────────────────────────────────
out = r"c:\Users\berek\OneDrive\Documents\bereket\SmartCampus_DSA_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
